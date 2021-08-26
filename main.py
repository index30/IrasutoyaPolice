# -*- coding: utf-8 -*-
"""Check the number of irasutoya images in ppt.

無償で使用して良いとされるいらすとやの画像の上限数を超過していないか検閲する

Example:
        $ python main.py

"""

from src import utils
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
import pickle
from pathlib import Path
import argparse
import itertools
from logging import basicConfig, DEBUG, ERROR, getLogger, \
    StreamHandler, WARNING

LIMIT = 20  # いらすとやの制限数
SITEMAP_URL = "https://www.irasutoya.com/sitemap.xml"  # スクレイピング用URL

# create logger
basicConfig(level=WARNING)
logger = getLogger(__name__)
logger.setLevel(WARNING)
# logger.addHandler(StreamHandler())


def check_irasutoya(args):
    """Check the number of irasutoya images in ppt.

    引数で入力したppt形式のファイルにいらすとやを使用した画像が何枚含まれているかを検出、
    規定を超えた量を使用していた場合警告を発出する関数。

    Args:
        args:   argparseによって受け取る引数

    """
    ppt_file = Presentation(args.ppt_filename)
    irasutoya_count = 0
    save_path = Path(Path(__file__).parent, 'assets', 'irasutoya.pickle')

    with open(save_path, 'rb') as irp:
        irstya_list = pickle.load(irp)
        irstya_keys = list(itertools.chain.from_iterable(zip(*irstya_list)))

    for slide in ppt_file.slides:
        for slide_shape in slide.shapes:
            if slide_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                ppt_img_name = slide_shape._pic.nvPicPr.cNvPr.get('descr')
                if ppt_img_name in irstya_keys:
                    irasutoya_count += 1

    logger.info("The number of your slide is {}.".format(irasutoya_count))
    if irasutoya_count > LIMIT:
        print("*******Caution*******\n"
              "If this slide is for commercial purpose, "
              "you must keep within {} irasutoya products.".format(LIMIT))
        print("*********************")


def scraping_irasutoya(args):
    result = utils.irasutoya_crawler(SITEMAP_URL)
    logger.debug("result size {}".format(len(result)))

    result2 = utils.find_img_names(result)
    logger.debug(result2[0])

    aseets_path = Path(Path(__file__).parent, 'assets')
    save_path = Path(aseets_path, 'irasutoya2.pickle')

    with open(save_path, 'wb') as isp:
        pickle.dump(result2, isp)

    with open(save_path, 'rb') as irp:
        result3 = pickle.load(irp)

    logger.debug(result3[0])


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Check irasutoya')
    subparsers = parser.add_subparsers()

    parser_check = subparsers.add_parser('check')
    parser_check.add_argument('ppt_filename')
    parser_check.set_defaults(handler=check_irasutoya)

    parser_scraping = subparsers.add_parser('scr')
    parser_scraping.set_defaults(handler=scraping_irasutoya)

    args = parser.parse_args()
    if hasattr(args, 'handler'):
        args.handler(args)
    else:
        parser.print_help()
